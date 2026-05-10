"""Document text extraction utilities for nanobot."""

import mimetypes
import os
from pathlib import Path
from typing import Any
from urllib.parse import urlparse

import httpx
from loguru import logger

from nanobot.utils.helpers import detect_image_mime

try:
    import magic
except ImportError:  # pragma: no cover - exercised only when optional dep is absent
    magic = None

# Supported file extensions for text extraction
SUPPORTED_EXTENSIONS: set[str] = {
    # Document formats
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    # Text formats
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".xml",
    ".html",
    ".htm",
    ".log",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    # Image formats (for future OCR support)
    ".png",
    ".jpg",
    ".jpeg",
    ".gif",
    ".webp",
}

_MAX_TEXT_LENGTH = 200_000


def extract_text(path: Path) -> str | None:
    """Extract text from a file.

    Args:
        path: Path to the file.

    Returns:
        Extracted text as string, None for unsupported types,
        or error string for failures.
    """
    if not isinstance(path, Path):
        path = Path(path)

    if not path.exists():
        return f"[error: file not found: {path}]"

    ext = path.suffix.lower()

    # Document formats -- each branch lazily imports its parser so that
    # startup does not pay the ~25 MB cost of loading openpyxl /
    # python-docx / python-pptx / pypdf up front (see issue #3422).
    if ext == ".pdf":
        return _extract_pdf(path)
    elif ext == ".docx":
        return _extract_docx(path)
    elif ext == ".xlsx":
        return _extract_xlsx(path)
    elif ext == ".pptx":
        return _extract_pptx(path)
    elif _is_text_extension(ext):
        return _extract_text_file(path)
    elif ext in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
        # Image files - for future OCR support
        return f"[image: {path.name}]"
    else:
        # Unsupported extension
        return None


def _extract_pdf(path: Path) -> str:
    """Extract text from PDF using pypdf."""
    try:
        from pypdf import PdfReader
    except ImportError:
        return "[error: pypdf not installed]"
    try:
        reader = PdfReader(path)
        pages: list[str] = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            pages.append(f"--- Page {i} ---\n{text}")
        return _truncate("\n\n".join(pages), _MAX_TEXT_LENGTH)
    except Exception as e:
        logger.error("Failed to extract PDF {}: {}", path, e)
        return f"[error: failed to extract PDF: {e!s}]"


def _extract_docx(path: Path) -> str:
    """Extract text from DOCX using python-docx."""
    try:
        from docx import Document as DocxDocument
    except ImportError:
        return "[error: python-docx not installed]"
    try:
        doc = DocxDocument(path)
        paragraphs: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
        return _truncate("\n\n".join(paragraphs), _MAX_TEXT_LENGTH)
    except Exception as e:
        logger.error("Failed to extract DOCX {}: {}", path, e)
        return f"[error: failed to extract DOCX: {e!s}]"


def _extract_xlsx(path: Path) -> str:
    """Extract text from XLSX using openpyxl."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "[error: openpyxl not installed]"
    try:
        wb = load_workbook(path, read_only=True, data_only=True)
        try:
            sheets: list[str] = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows: list[str] = []
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join(
                        str(cell) if cell is not None else "" for cell in row
                    )
                    if row_text.strip():
                        rows.append(row_text)
                if rows:
                    sheets.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
            return _truncate("\n\n".join(sheets), _MAX_TEXT_LENGTH)
        finally:
            wb.close()
    except Exception as e:
        logger.error("Failed to extract XLSX {}: {}", path, e)
        return f"[error: failed to extract XLSX: {e!s}]"


def _extract_pptx(path: Path) -> str:
    """Extract text from PPTX using python-pptx."""
    try:
        from pptx import Presentation as PptxPresentation
    except ImportError:
        return "[error: python-pptx not installed]"
    try:
        prs = PptxPresentation(path)
        slides: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text: list[str] = []
            for shape in slide.shapes:
                _collect_pptx_shape_text(shape, slide_text)
            if slide_text:
                slides.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))
        return _truncate("\n\n".join(slides), _MAX_TEXT_LENGTH)
    except Exception as e:
        logger.error("Failed to extract PPTX {}: {}", path, e)
        return f"[error: failed to extract PPTX: {e!s}]"


def _collect_pptx_shape_text(shape, out: list[str]) -> None:
    """Collect text from a PPTX shape, recursing into groups and tables.

    Groups have ``has_text_frame=False`` and must be walked via ``.shapes``;
    tables are GraphicFrame objects whose cell text lives under ``.table``.
    """
    sub_shapes = getattr(shape, "shapes", None)
    if sub_shapes is not None:
        for sub in sub_shapes:
            _collect_pptx_shape_text(sub, out)
        return

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            line = "\t".join(cell for cell in cells if cell)
            if line:
                out.append(line)
        return

    text = getattr(shape, "text", "")
    if text:
        out.append(text)


def _extract_text_file(path: Path) -> str:
    """Extract text from a plain text file."""
    try:
        # Try UTF-8 first, then latin-1 fallback
        try:
            content = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            content = path.read_text(encoding="latin-1")
        return _truncate(content, _MAX_TEXT_LENGTH)
    except Exception as e:
        logger.error("Failed to read text file {}: {}", path, e)
        return f"[error: failed to read file: {e!s}]"


def _truncate(text: str, max_length: int) -> str:
    """Truncate text with a suffix indicating truncation."""
    if len(text) <= max_length:
        return text
    return text[:max_length] + f"... (truncated, {len(text)} chars total)"


def _is_text_extension(ext: str) -> bool:
    """Check if extension is a text format."""
    return ext in {
        ".txt",
        ".md",
        ".csv",
        ".json",
        ".xml",
        ".html",
        ".htm",
        ".log",
        ".yaml",
        ".yml",
        ".toml",
        ".ini",
        ".cfg",
    }


# ---------------------------------------------------------------------------
# High-level helper: normalize attachments into multimodal URLs + text refs
# ---------------------------------------------------------------------------

_MAX_EXTRACT_FILE_SIZE = 50 * 1024 * 1024  # 50 MB

_VISUAL_MIME_PREFIXES = ("image/", "video/")
_VISUAL_REF_PREFIXES = {
    "image_url": "image_url:",
    "video_url": "video_url:",
}


def _is_remote_attachment_ref(ref: str) -> bool:
    parsed = urlparse(ref)
    return parsed.scheme in {"http", "https"}


def _read_header(path: Path, size: int = 64) -> bytes:
    with path.open("rb") as f:
        return f.read(size)


def _detect_remote_mime_by_content(
    url: str,
    *,
    timeout: float = 10,
    header_size: int = 4096,
) -> str | None:
    headers = {
        "User-Agent": "Mozilla/5.0",
        "Accept": "*/*",
        "Range": f"bytes=0-{header_size - 1}",
    }

    try:
        with httpx.Client(follow_redirects=True, timeout=timeout) as client:
            resp = client.get(url, headers=headers)
            resp.raise_for_status()

            header = resp.content[:header_size]
            if not header:
                return None

            if magic is None:
                return None
            return magic.from_buffer(header, mime=True)

    except httpx.HTTPError:
        return None


def _detect_attachment_mime(ref: str, path: Path | None = None) -> str | None:
    if path is not None and path.is_file():
        try:
            header = _read_header(path)
        except OSError:
            header = b""
        detected = detect_image_mime(header)
        if detected:
            return detected
        return mimetypes.guess_type(path.name, strict=False)[0]

    return _detect_remote_mime_by_content(ref)


def _is_visual_mime(mime: str | None) -> bool:
    return bool(mime and mime.startswith(_VISUAL_MIME_PREFIXES))


def _visual_block_type_for_mime(mime: str) -> str | None:
    if mime.startswith("image/"):
        return "image_url"
    if mime.startswith("video/"):
        return "video_url"
    return None


def _is_vllm_provider_backend(
    provider: Any | None = None,
    metadata: dict[str, Any] | None = None,
) -> bool:
    meta = metadata or {}
    provider_name = str(
        meta.get("provider_backend")
        or meta.get("provider_name")
        or meta.get("provider")
        or ""
    ).strip().lower()
    if provider_name == "vllm":
        return True

    spec = getattr(provider, "_spec", None)
    spec_name = str(getattr(spec, "name", "") or "").strip().lower()
    return spec_name == "vllm"


def encode_visual_media_ref(block_type: str, url: str) -> str:
    prefix = _VISUAL_REF_PREFIXES.get(block_type)
    if not prefix:
        raise ValueError(f"unsupported visual block type: {block_type}")
    return f"{prefix}{url}"


def _attachment_reference_line(mime: str | None, ref: str) -> str:
    file_type = mime or "application/octet-stream"
    return f"“{file_type}”：{ref}"


def _container_up_attachment_upload_url() -> str:
    return str(os.environ.get("CONTAINER_UP_ATTACHMENT_UPLOAD_URL") or "").strip()


def _upload_local_attachment_ref(
    path: Path,
    *,
    metadata: dict[str, Any] | None = None,
) -> str | None:
    meta = dict(metadata or {})
    frontend_id = str(meta.get("frontend_id") or "").strip()
    user_id = str(meta.get("usr_id") or meta.get("user_id") or "user").strip() or "user"
    if not frontend_id:
        logger.warning(
            "Skipping local attachment upload for {}: missing frontend_id", path
        )
        return None

    url = _container_up_attachment_upload_url()
    if not url:
        logger.warning(
            "Skipping local attachment upload for {}: missing CONTAINER_UP_ATTACHMENT_UPLOAD_URL",
            path,
        )
        return None
    try:
        with httpx.Client(
            timeout=float(
                os.environ.get("CONTAINER_UP_ATTACHMENT_UPLOAD_TIMEOUT_SECONDS", "30")
            ),
        ) as client:
            response = client.post(
                url,
                json={
                    "frontend_id": frontend_id,
                    "user_id": user_id,
                    "local_path": str(path),
                },
            )
            response.raise_for_status()
            payload = response.json()
    except Exception as exc:
        logger.warning(
            "Failed to upload local attachment {} via container-up: {}", path, exc
        )
        return None
    uploaded = str(payload.get("url") or "").strip()
    if not uploaded:
        logger.warning("container-up upload returned no URL for {}", path)
        return None
    return uploaded


def extract_documents(
    text: str,
    media_paths: list[str],
    *,
    max_file_size: int = _MAX_EXTRACT_FILE_SIZE,
    metadata: dict[str, Any] | None = None,
    provider: Any | None = None,
) -> tuple[str, list[str]]:
    """Normalize inbound attachments into vision URLs and text references.

    Visual attachments (images/videos) are returned in the media list as
    ``image_url:<url>`` / ``video_url:<url>`` records so downstream prompt
    builders do not need to re-detect MIME types.
    Ordinary files are appended to the text content in the form
    ``“<mime>”：<url>``.

    Local attachments are uploaded to the configured MinIO ``attachments``
    bucket first; remote HTTP(S) URLs are consumed directly.
    """
    supports_visual_url = _is_vllm_provider_backend(provider=provider, metadata=metadata)
    visual_refs: list[str] = []
    attachment_refs: list[str] = []

    for ref in media_paths:
        ref_str = str(ref or "").strip()
        if not ref_str:
            continue

        local_path = Path(ref_str).expanduser()
        resolved_ref = ref_str
        mime: str | None = None

        if local_path.is_file():
            try:
                size = local_path.stat().st_size
            except OSError:
                continue
            if size > max_file_size:
                logger.warning(
                    "Skipping oversized file for extraction: {} ({:.1f} MB > {} MB limit)",
                    local_path.name,
                    size / (1024 * 1024),
                    max_file_size // (1024 * 1024),
                )
                continue
            mime = _detect_attachment_mime(ref_str, local_path)
            uploaded_ref = _upload_local_attachment_ref(local_path, metadata=metadata)
            if not uploaded_ref:
                continue
            resolved_ref = uploaded_ref
        else:
            mime = _detect_attachment_mime(ref_str)
            if not _is_remote_attachment_ref(ref_str):
                logger.warning(
                    "Skipping unsupported non-local attachment reference: {}", ref_str
                )
                continue

        if supports_visual_url and _is_visual_mime(mime):
            block_type = _visual_block_type_for_mime(mime)
            if not block_type:
                continue
            visual_refs.append(encode_visual_media_ref(block_type, resolved_ref))
        else:
            attachment_refs.append(_attachment_reference_line(mime, resolved_ref))

    if attachment_refs:
        text = (
            text + "\n\n" + "\n\n".join(attachment_refs)
            if text
            else "\n\n".join(attachment_refs)
        )

    return text, visual_refs
